#!/usr/bin/env python3
"""Generate Anitrack_Defense.pptx from 演讲大纲.md structure."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
FIGURES = ROOT.parent / "anitrack-visuals" / "figures"
OUT = ROOT / "slides" / "Anitrack_Defense.pptx"

# Theme
BG = RGBColor(0x0F, 0x17, 0x2A)
TITLE_COLOR = RGBColor(0xF8, 0xFA, 0xFC)
BODY_COLOR = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
TABLE_HEADER = RGBColor(0x1E, 0x40, 0xAF)
TABLE_ROW_ALT = RGBColor(0x1E, 0x29, 0x3B)


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_box(slide, text, top=Inches(0.35), height=Inches(0.9), size=32):
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(9.0), height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.LEFT
    return box


def add_bullets(slide, lines, top=Inches(1.35), left=Inches(0.55), width=Inches(8.8), size=18):
    box = slide.shapes.add_textbox(left, top, width, Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = BODY_COLOR
        p.space_after = Pt(10)
    return box


def add_subtitle(slide, text, top=Inches(1.05), size=14):
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(9.0), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.italic = True
    p.font.color.rgb = MUTED


def add_table(slide, headers, rows, top=Inches(1.3), col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    width = Inches(9.0)
    height = Inches(0.42 * n_rows)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = TITLE_COLOR
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = BODY_COLOR
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_ALT
    return table_shape


def add_image_if_exists(slide, name, left, top, width, height=None):
    path = FIGURES / name
    if not path.exists():
        box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.text = f"[Missing: {name}]"
        p.font.color.rgb = MUTED
        return None
    if height:
        return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    return slide.shapes.add_picture(str(path), left, top, width=width)


def add_placeholder(slide, text, top=Inches(2.2), left=Inches(0.55), height=Inches(2.8), width=Inches(4.2)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    shape.line.color.rgb = ACCENT
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.color.rgb = MUTED


def build():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9
    blank = prs.slide_layouts[6]

    # 1 — Cover
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    lines = [
        "Anitrack",
        "Personal Anime Watchlist & Analytics",
        "Web Technologies I · HSBI · May 2026",
        "[Name 1] · [Name 2] · [Name 3]",
    ]
    box = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.2))
    tf = box.text_frame
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        if i == 0:
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = ACCENT
        elif i == 1:
            p.font.size = Pt(22)
            p.font.color.rgb = TITLE_COLOR
        else:
            p.font.size = Pt(16)
            p.font.color.rgb = MUTED
        p.space_after = Pt(14)

    # 2 — Team
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_box(s, "Team & responsibilities")
    add_table(
        s,
        ["Member", "Main focus"],
        [
            ["[Name A]", "Frontend — Next.js, UI, i18n, responsive"],
            ["[Name B]", "Backend — NestJS, MongoDB, Bee"],
            ["[Name C]", "API contract, tests, live demo"],
        ],
        col_widths=[Inches(2.2), Inches(6.8)],
    )

    # 3 — Problem & goals
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_box(s, "Problem & goals")
    add_bullets(
        s,
        [
            "Track anime: Planned → Watching → Completed",
            "Monthly heatmap (GitHub-style)",
            "Seasonal timetable & dashboard picks (local mirror)",
            "API-first full-stack for course requirements",
        ],
    )

    # 4 — Course requirements
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_box(s, "Course requirements ↔ solution")
    add_subtitle(
        s,
        "Content requirements: implemented. Organizational requirements: fulfilled in this presentation.",
        top=Inches(1.02),
        size=12,
    )
    add_table(
        s,
        ["Requirement", "Anitrack"],
        [
            ["NestJS backend", "anitrack-backend :3001"],
            ["CRUD HTTP API", "/api/anime"],
            ["API first", "swagger.json, /api-docs"],
            ["Tests", "Jest, Vitest, contract-validator"],
            ["Thin frontend", "Logic in backend"],
            ["Responsive", "Mobile + desktop"],
        ],
        top=Inches(1.45),
        col_widths=[Inches(3.2), Inches(5.8)],
    )

    # 5 — Tech stack
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_box(s, "Technology stack")
    add_subtitle(s, "Frontend · Backend · MongoDB · Jikan · Bangumi")
    add_image_if_exists(s, "tech-stack.png", Inches(0.4), Inches(1.35), Inches(9.2))

    # 6 — Architecture
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_box(s, "System architecture")
    add_subtitle(s, "Bee mirrors Jikan/Bangumi → MongoDB; client → NestJS API only")
    add_image_if_exists(s, "architecture.png", Inches(0.35), Inches(1.32), Inches(9.3))

    # 7 — API-first
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_box(s, "API-first & OpenAPI")
    add_bullets(
        s,
        [
            "Swagger UI: http://localhost:3001/api-docs",
            "Contract: anitrack-backend/swagger.json",
            "Same API for UI + script clients (anitrack-tester)",
            "Errors: { error: { code, message, details } }",
        ],
        top=Inches(1.25),
        size=17,
    )
    add_placeholder(
        s,
        "Insert screenshot:\nlocalhost:3001/api-docs\n(Project_Intro/screenshots/)",
        top=Inches(2.85),
        height=Inches(2.2),
        width=Inches(8.9),
    )

    # 8 — User flows
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_box(s, "User flows")
    add_subtitle(s, "Dashboard · Timetable · Library · Profile")
    add_image_if_exists(s, "user-flow.png", Inches(0.35), Inches(1.28), Inches(9.3))

    # 9 — Data flow
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_box(s, "Data flow & caching")
    add_image_if_exists(
        s, "data-flow.png", Inches(4.6), Inches(1.25), Inches(5.0), height=Inches(3.9)
    )
    add_bullets(
        s,
        [
            "Cache-aside AnimeMeta by malId",
            "User progress → AnimeEntry",
            "Bee → AnimeMirror for timetable / seasonal",
        ],
        top=Inches(1.35),
        left=Inches(0.5),
        width=Inches(4.0),
        size=17,
    )

    # 10 — Backend logic
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_box(s, "Backend business logic")
    add_bullets(
        s,
        [
            "State machine → illegal status = HTTP 409",
            "Heatmap aggregation on server (intensity 0–4)",
            "Bee: 65s / 3 requests, Bangumi mapping",
            "No business rules in React frontend",
        ],
    )

    # 11 — Testing
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_box(s, "Testing strategy")
    add_bullets(
        s,
        [
            "Unit: heatmap mapping (Vitest)",
            "Integration: API + MongoDB (Jest, Vitest)",
            "Contract: OpenAPI vs runtime (anitrack-tester)",
        ],
    )

    # 12 — Responsive
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_box(s, "Responsive UI")
    add_bullets(
        s,
        [
            "Hamburger navigation · date strip ±2 weeks · heatmap scroll",
        ],
        top=Inches(1.15),
        size=16,
    )
    add_placeholder(
        s,
        "Mobile / narrow\n(insert screenshot)",
        top=Inches(1.75),
        left=Inches(0.55),
        height=Inches(3.5),
    )
    add_placeholder(
        s,
        "Desktop / wide\n(insert screenshot)",
        top=Inches(1.75),
        left=Inches(5.0),
        height=Inches(3.5),
    )

    # 13 — Demo plan
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_box(s, "Live demo plan")
    add_bullets(
        s,
        [
            "1. Dashboard — watching & seasonal picks",
            "2. Library — search → add to watchlist",
            "3. Profile — heatmap → month activity",
            "4. (Optional) Timetable — pick date",
            "",
            "Before demo: backend :3001, frontend :3000, UI = English",
        ],
    )

    # 14 — Live demo (transition slide)
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    box = s.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(8.0), Inches(2.0))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Live Demo"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.text = "http://localhost:3000"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(22)
    p2.font.color.rgb = TITLE_COLOR
    p3 = tf.add_paragraph()
    p3.text = "(Switch to browser — no further slides during demo)"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(14)
    p3.font.color.rgb = MUTED

    # 15 — Limitations
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_box(s, "Limitations & outlook")
    add_bullets(
        s,
        [
            "Timetable air times sometimes TBD (upstream data)",
            "Single user (TEMP_USER_ID) — no auth yet",
            "Future: multi-user, recommendations",
            "Course core scope is complete",
        ],
    )

    # 16 — Thank you
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    box = s.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(8.0), Inches(1.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank you for your attention."
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p2 = tf.add_paragraph()
    p2.text = "Questions?"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(28)
    p2.font.color.rgb = ACCENT

    prs.save(str(OUT))
    print(f"Saved: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
